#!/usr/bin/env python3
"""
Convert HTML Reveal.js presentations to PowerPoint format.
"""

import os
import re
from html.parser import HTMLParser
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ACM Brand Colors
ACM_GREEN = RGBColor(0x00, 0x68, 0x37)
ACM_GOLD = RGBColor(0xD4, 0xAF, 0x37)
ACM_DARK = RGBColor(0x1a, 0x1a, 0x2e)
ACM_PRIMARY = RGBColor(0x1a, 0x5f, 0x7a)


class SlideExtractor(HTMLParser):
    """Extract slide content from HTML."""

    def __init__(self):
        super().__init__()
        self.slides = []
        self.current_slide = None
        self.current_text = ""
        self.in_section = False
        self.in_h1 = False
        self.in_h2 = False
        self.in_h3 = False
        self.in_li = False
        self.in_td = False
        self.in_th = False
        self.in_pre = False
        self.in_code = False
        self.tag_stack = []
        self.current_list_items = []
        self.current_table = []
        self.current_row = []

    def handle_starttag(self, tag, attrs):
        self.tag_stack.append(tag)
        attrs_dict = dict(attrs)

        if tag == "section":
            self.in_section = True
            if self.current_slide:
                self.slides.append(self.current_slide)
            self.current_slide = {
                "title": "",
                "subtitle": "",
                "content": [],
                "bullets": [],
                "table": [],
                "code": "",
                "is_title_slide": "title-slide" in attrs_dict.get("class", "")
            }
        elif tag == "h1":
            self.in_h1 = True
            self.current_text = ""
        elif tag == "h2":
            self.in_h2 = True
            self.current_text = ""
        elif tag == "h3":
            self.in_h3 = True
            self.current_text = ""
        elif tag == "li":
            self.in_li = True
            self.current_text = ""
        elif tag == "td":
            self.in_td = True
            self.current_text = ""
        elif tag == "th":
            self.in_th = True
            self.current_text = ""
        elif tag == "tr":
            self.current_row = []
        elif tag == "pre":
            self.in_pre = True
            self.current_text = ""
        elif tag == "code":
            self.in_code = True

    def handle_endtag(self, tag):
        if self.tag_stack and self.tag_stack[-1] == tag:
            self.tag_stack.pop()

        if tag == "section":
            self.in_section = False
        elif tag == "h1" and self.current_slide:
            self.in_h1 = False
            self.current_slide["title"] = self.clean_text(self.current_text)
        elif tag == "h2" and self.current_slide:
            self.in_h2 = False
            if not self.current_slide["title"]:
                self.current_slide["title"] = self.clean_text(self.current_text)
            else:
                self.current_slide["subtitle"] = self.clean_text(self.current_text)
        elif tag == "h3" and self.current_slide:
            self.in_h3 = False
            self.current_slide["content"].append(("h3", self.clean_text(self.current_text)))
        elif tag == "li" and self.current_slide:
            self.in_li = False
            text = self.clean_text(self.current_text)
            if text:
                self.current_slide["bullets"].append(text)
        elif tag == "td" or tag == "th":
            self.in_td = False
            self.in_th = False
            self.current_row.append(self.clean_text(self.current_text))
        elif tag == "tr" and self.current_slide and self.current_row:
            self.current_slide["table"].append(self.current_row)
            self.current_row = []
        elif tag == "pre" and self.current_slide:
            self.in_pre = False
            self.current_slide["code"] = self.current_text
        elif tag == "code":
            self.in_code = False

    def handle_data(self, data):
        if self.in_h1 or self.in_h2 or self.in_h3 or self.in_li or self.in_td or self.in_th:
            self.current_text += data
        elif self.in_pre:
            self.current_text += data

    def clean_text(self, text):
        """Clean and normalize text."""
        text = re.sub(r'\s+', ' ', text)
        text = text.strip()
        return text

    def get_slides(self):
        if self.current_slide:
            self.slides.append(self.current_slide)
        return self.slides


def create_title_slide(prs, slide_data):
    """Create a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add green header bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACM_GREEN
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = slide_data.get("title", "")
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = ACM_GREEN
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    if slide_data.get("subtitle"):
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = slide_data.get("subtitle", "")
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = ACM_DARK
        subtitle_para.alignment = PP_ALIGN.CENTER

    return slide


def create_content_slide(prs, slide_data):
    """Create a content slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add green header bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACM_GREEN
    shape.line.fill.background()

    # Title
    if slide_data.get("title"):
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data["title"]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = ACM_GREEN

    # Subtitle
    y_pos = 1.3
    if slide_data.get("subtitle"):
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.5))
        sub_frame = sub_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = slide_data["subtitle"]
        sub_para.font.size = Pt(20)
        sub_para.font.color.rgb = ACM_DARK
        y_pos += 0.6

    # Bullets
    if slide_data.get("bullets"):
        bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(4))
        bullet_frame = bullet_box.text_frame
        bullet_frame.word_wrap = True

        for i, bullet in enumerate(slide_data["bullets"][:10]):  # Max 10 bullets
            if i == 0:
                para = bullet_frame.paragraphs[0]
            else:
                para = bullet_frame.add_paragraph()
            para.text = "• " + bullet
            para.font.size = Pt(18)
            para.font.color.rgb = ACM_DARK
            para.space_after = Pt(8)

    # Table
    if slide_data.get("table") and len(slide_data["table"]) > 1:
        table_data = slide_data["table"]
        rows = min(len(table_data), 8)
        cols = min(len(table_data[0]), 5) if table_data else 0

        if cols > 0:
            table_top = Inches(y_pos if not slide_data.get("bullets") else 4.5)
            table = slide.shapes.add_table(rows, cols, Inches(0.5), table_top, Inches(9), Inches(2.5)).table

            for i, row_data in enumerate(table_data[:rows]):
                for j, cell_text in enumerate(row_data[:cols]):
                    cell = table.cell(i, j)
                    cell.text = str(cell_text)
                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(12)

                    if i == 0:  # Header row
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = ACM_GREEN
                        para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        para.font.bold = True

    return slide


def convert_html_to_pptx(html_path, output_path, title="Presentation"):
    """Convert an HTML Reveal.js presentation to PowerPoint."""

    # Read HTML
    with open(html_path, 'r', encoding='utf-8') as f:
        html_content = f.read()

    # Parse slides
    parser = SlideExtractor()
    parser.feed(html_content)
    slides = parser.get_slides()

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create slides
    for slide_data in slides:
        if not slide_data.get("title") and not slide_data.get("bullets"):
            continue

        if slide_data.get("is_title_slide"):
            create_title_slide(prs, slide_data)
        else:
            create_content_slide(prs, slide_data)

    # Save
    prs.save(output_path)
    print(f"Created: {output_path}")


def main():
    base_dir = "/Users/AbiolaOgunsakin1/Anti_Call-Masking/anti-call-masking/docs"

    # Training presentations
    training_presentations = [
        ("training/presentations/01-system-overview.html", "training/presentations/exports/pptx/01-system-overview.pptx"),
        ("training/presentations/02-admin-training.html", "training/presentations/exports/pptx/02-admin-training.pptx"),
        ("training/presentations/03-soc-analyst-training.html", "training/presentations/exports/pptx/03-soc-analyst-training.pptx"),
        ("training/presentations/04-developer-training.html", "training/presentations/exports/pptx/04-developer-training.pptx"),
    ]

    # NCC presentations
    ncc_presentations = [
        ("presentations/NCC_EXECUTIVE_SUMMARY.html", "presentations/exports/pptx/NCC_EXECUTIVE_SUMMARY.pptx"),
        ("presentations/NCC_CAPABILITIES_DEMO.html", "presentations/exports/pptx/NCC_CAPABILITIES_DEMO.pptx"),
    ]

    all_presentations = training_presentations + ncc_presentations

    for html_file, pptx_file in all_presentations:
        html_path = os.path.join(base_dir, html_file)
        pptx_path = os.path.join(base_dir, pptx_file)

        if os.path.exists(html_path):
            try:
                convert_html_to_pptx(html_path, pptx_path)
            except Exception as e:
                print(f"Error converting {html_file}: {e}")
        else:
            print(f"Not found: {html_path}")


if __name__ == "__main__":
    main()
